import os
import tempfile
import hashlib
import uuid
from dotenv import load_dotenv
from langchain.embeddings import SentenceTransformerEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from langchain_community.vectorstores import Chroma
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
import chromadb
import requests
from bs4 import BeautifulSoup
from openai import OpenAI
from docx import Document
from pptx import Presentation

load_dotenv()

os.environ['VECTOR_DB'] = './db'

vector_db_path = os.getenv('VECTOR_DB')
print(f"VECTOR_DB path: {vector_db_path}")

if not vector_db_path or vector_db_path.startswith('http'):
    raise ValueError("The VECTOR_DB environment variable is not correctly set. It should be a valid directory path.")

if not os.path.isdir(vector_db_path):
    os.makedirs(vector_db_path)

api_key = "sk-taral-upwork-VqgjaTN5HZGDFmE9sIAAT3BlbkFJ8AkY3JurXEUjXh1cfnnj"
client = OpenAI(api_key=api_key)

class CustomSentenceTransformerEmbeddings:
    def __init__(self, model_name):
        self.model = SentenceTransformer(model_name)

    def embed_documents(self, texts):
        return self.model.encode(texts).tolist() 

    def embed_query(self, text):
        return self.embed_documents([text])[0]

    def __call__(self, input):
        return self.embed_documents(input)

embedding_function = CustomSentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")

def connect_db(collection_name):
    vector_db_path = os.getenv('VECTOR_DB')
    if not vector_db_path:
        raise ValueError("The VECTOR_DB environment variable is not set.")
    
    print(f"Connecting to ChromaDB at {vector_db_path}...")
    try:
        client = chromadb.PersistentClient(path=vector_db_path)
        collection = client.get_or_create_collection(
            collection_name,
            embedding_function=embedding_function,
            metadata={"hnsw:space": "cosine"}
        )
        print("Connected to ChromaDB successfully.")
        return collection
    except Exception as e:
        print(f"Failed to connect to ChromaDB: {e}")
        raise

def is_valid_chunk(chunk_text):
    if len(chunk_text.strip()) == 0:
        return False
    if "<" in chunk_text or ">" in chunk_text:
        return False
    if "{" in chunk_text or "}" in chunk_text:
        return False
    if any(keyword in chunk_text.lower() for keyword in ["metadata", "object", "instanceid", "rdf", "xmlns"]):
        return False
    return True

def store_chunks(text, prefix):
    chunks = []
    chunk_size = 200
    text_chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    for idx, chunk_text in enumerate(text_chunks):
        if is_valid_chunk(chunk_text):
            chunk = {
                'id': f"{prefix}_{idx}",
                'metadata': {'text': chunk_text}
            }
            chunks.append(chunk)
    return chunks

def fetch_company_info_from_link(link):
    response = requests.get(link)
    soup = BeautifulSoup(response.text, 'html.parser')
    return soup.get_text()

def extract_text_from_file(filepath):
    if filepath.endswith('.pdf'):
        reader = PdfReader(filepath)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
    elif filepath.endswith('.docx'):
        from docx import Document
        doc = Document(filepath)
        text = ' '.join([para.text for para in doc.paragraphs])
    elif filepath.endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif filepath.endswith('.txt'):
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
    else:
        text = ''
    return text

def read_document(file_path, user):
    filename = os.path.basename(file_path)
    with open(file_path, 'rb') as f:
        data = f.read()

    doc_hash = hashlib.sha256(data).hexdigest()
    collection = connect_db("document_chunks")
    check_duplicates = collection.get(
        where={
            '$and': [
                {'user': {'$in': ['PUBLIC', user]}},
                {'doc_hash': doc_hash}
            ]
        }
    )['ids']
    if check_duplicates:
        return ('Dupe', 'Dupe', 'Dupe')

    text = extract_text_from_file(file_path)
    return filename, text, doc_hash

def extract_text_from_pdf(filepath):
    reader = PdfReader(filepath)
    text = ''
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(filepath):
    import docx
    doc = docx.Document(filepath)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_xlsx(filepath):
    import pandas as pd
    df = pd.read_excel(filepath)
    text = df.to_string()
    return text

def extract_text_from_txt(filepath):
    with open(filepath, 'r', encoding='utf-8') as file:
        text = file.read()
    return text

def process_and_store_files(transcript, company_info, company_info_link):
    all_chunks = []
    if transcript:
        transcript_chunks = store_chunks(transcript, "transcript")
        all_chunks.extend(transcript_chunks)
    if company_info:
        company_info_chunks = store_chunks(company_info, "company_info")
        all_chunks.extend(company_info_chunks)
    if company_info_link:
        company_info_text = fetch_company_info_from_link(company_info_link)
        company_info_link_chunks = store_chunks(company_info_text, "company_info_link")
        all_chunks.extend(company_info_link_chunks)

    print("All Chunks Created:")
    for chunk in all_chunks:
        print(f"Chunk ID: {chunk['id']}, Text: {chunk['metadata']['text']}")

    if all_chunks:
        collection = connect_db("document_chunks")
        ids = [doc["id"] for doc in all_chunks]
        embeddings = [embedding_function.embed_documents([doc["metadata"]["text"]])[0] for doc in all_chunks]
        metadatas = [doc["metadata"] for doc in all_chunks]
        collection.add(ids=ids, embeddings=embeddings, metadatas=metadatas)

def search_relevant_chunks(query, top_k=5):
    query_embedding = embedding_function.embed_query(query)
    collection = connect_db("document_chunks")
    results = collection.query(query_embeddings=[query_embedding], n_results=top_k)
    return results["documents"]

def get_openai_response(combined_text, query):
    prompt = f"""
    You are a helpful assistant. Please answer the following question based only on the provided context.
    
    Context:
    {combined_text}
    
    Question:
    {query}
    
    Answer:
    """

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=150,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    answer = response.choices[0].message.content.strip()
    return answer

def retrieve_and_create_chain(query):
    query_embedding = embedding_function.embed_query(query)
    collection = connect_db("document_chunks")
    results = collection.query(query_embeddings=[query_embedding], n_results=5)
    docs = results["documents"]
    
    DEFAULT_DOCUMENT_PROMPT = PromptTemplate.from_template(template="{page_content}")
    def _combine_documents(docs, document_prompt=DEFAULT_DOCUMENT_PROMPT, document_separator="\n\n"):
        doc_strings = [doc['metadata']['text'] for doc in docs]
        return document_separator.join(doc_strings)

    context = _combine_documents(docs)
    template = """You are a helpful assistant. Given the following context, please provide a detailed and specific answer to the question provided.
    Context:
    {context}
    
    Question:
    {question}
    
    Answer:"""
    PROMPT = PromptTemplate.from_template(template)

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a helpful assistant. Given the following context, please provide a detailed and specific answer to the question provided."},
            {"role": "user", "content": f"Context:\n{context}\n\nQuestion:\n{query}\n\nAnswer:"}
        ],
        temperature=0.7,
        max_tokens=150,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    answer = response.choices[0].message.content.strip()
    return answer
